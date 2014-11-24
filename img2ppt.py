# -*- coding: utf-8 -*-
"""
Created on Mon Nov 24 19:38:32 2014

@author: tobie

Simple function to take a list of paths to filenames and create a 
powerpoint presentation with one image per slide

"""

from pptx import Presentation
from pptx.util import Inches
import glob

    
def img2ppt( pics, outfile ):
    
    """ takes a list of pictures and creates a slideshow in powerpoint
        pics    = list of pics
        outfile = filename.pptx
    """
    pres = Presentation()
    for pic in pics:
        slide = pres.slides.add_slide( pres.slide_layouts[1] )
        slide.shapes.add_picture( pic, Inches(0.6), Inches(2))
    pres.save(outfile)
    
def test():
    """need to run in ipython"""
    pics = glob.glob("*.png")
    filename = 'test.pptx'
    img2ppt( pics, filename)
